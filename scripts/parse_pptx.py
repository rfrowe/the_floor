#!/usr/bin/env python3
"""
PPTX Parser for "The Floor" Game Show Application

This script parses a PowerPoint (PPTX) file exported from Google Slides and
extracts:
- Slide images (as base64)
- Speaker notes (as answers)
- Censorship boxes (rectangles with specific properties)

Output is JSON format compatible with the TypeScript data models.

Usage:
    python scripts/parse_pptx.py input.pptx output.json --contestant "John Doe" --category "Movies"

Requirements:
    pip install python-pptx pillow
"""

import argparse
import base64
import json
import sys
from dataclasses import dataclass, asdict
from io import BytesIO
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
    from PIL import Image
except ImportError as e:
    print(f"Error: Missing required library: {e}", file=sys.stderr)
    print("Install with: pip install python-pptx pillow", file=sys.stderr)
    sys.exit(1)


# Type definitions that mirror TypeScript types in src/types/
# IMPORTANT: Keep these in sync with TypeScript types!
# If you update these, update src/types/slide.ts and src/types/contestant.ts


@dataclass
class CensorBox:
    """
    Represents a censorship box overlay on a slide image.
    Position and size are specified as percentages (0-100) relative to the image dimensions.

    SYNC WITH: src/types/slide.ts - CensorBox interface
    """
    x: float  # X position as percentage (0-100) from left edge
    y: float  # Y position as percentage (0-100) from top edge
    width: float  # Width as percentage (0-100) of image width
    height: float  # Height as percentage (0-100) of image height
    color: str  # Color of the censor box (hex format)


@dataclass
class Slide:
    """
    Represents a single slide in a category, containing an image,
    the correct answer, and optional censorship boxes.

    SYNC WITH: src/types/slide.ts - Slide interface
    """
    imageUrl: str  # Image data as base64 data URL
    answer: str  # The correct answer for this slide (from speaker notes)
    censorBoxes: list[CensorBox]  # Censorship boxes to overlay on the image


@dataclass
class Category:
    """
    Represents a category (topic) with a collection of slides for gameplay.

    SYNC WITH: src/types/contestant.ts - Category interface
    """
    name: str
    slides: list[Slide]


@dataclass
class ParsedData:
    """
    Output format for the parsed PPTX data.
    Contains category data and optional metadata about contestant.
    """
    category: Category
    metadata: dict[str, str] | None = None


def emu_to_percentage(emu_value: int, slide_dimension_emu: int) -> float:
    """Convert EMU (English Metric Units) to percentage of slide dimension."""
    return (emu_value / slide_dimension_emu) * 100


def extract_image_as_base64(slide, slide_index: int) -> str | None:
    """
    Extract the main image from a slide and convert to base64.
    Applies any cropping that was set in the PPTX.
    """
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            image = shape.image
            image_bytes = image.blob

            # Get crop information
            crop_left, crop_top, crop_right, crop_bottom = get_image_crop_info(shape)

            # Convert to PNG with white background if needed
            try:
                img = Image.open(BytesIO(image_bytes))

                # Apply cropping if specified
                if any([crop_left, crop_top, crop_right, crop_bottom]):
                    width, height = img.size

                    # Calculate crop box in pixels
                    left_px = int(width * crop_left)
                    top_px = int(height * crop_top)
                    right_px = int(width * (1 - crop_right))
                    bottom_px = int(height * (1 - crop_bottom))

                    # Crop the image
                    img = img.crop((left_px, top_px, right_px, bottom_px))
                    print(f"  Applied crop to image: {crop_left*100:.1f}%/{crop_top*100:.1f}%/{crop_right*100:.1f}%/{crop_bottom*100:.1f}%, new size: {img.size}", file=sys.stderr)

                # If image has transparency, add white background
                if img.mode in ('RGBA', 'LA') or (img.mode == 'P' and 'transparency' in img.info):
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    if img.mode == 'P':
                        img = img.convert('RGBA')
                    background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                    img = background

                # Convert to base64
                buffer = BytesIO()
                img.save(buffer, format='PNG')
                img_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
                return f"data:image/png;base64,{img_base64}"
            except Exception as e:
                print(f"Warning: Failed to process image on slide {slide_index + 1}: {e}", file=sys.stderr)
                # Fallback: return original image as base64
                img_base64 = base64.b64encode(image_bytes).decode('utf-8')
                ext = image.ext or 'png'
                return f"data:image/{ext};base64,{img_base64}"

    return None


def extract_speaker_notes(slide) -> str:
    """Extract speaker notes from a slide."""
    if not slide.has_notes_slide:
        return ""

    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    return text_frame.text.strip()


def find_image_shape(slide):
    """Find the main image shape on a slide."""
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            return shape
    return None


def get_image_crop_info(image_shape):
    """
    Get crop information for an image shape.
    Returns (crop_left, crop_top, crop_right, crop_bottom) as fractions (0.0 to 1.0).
    """
    crop_left = image_shape.crop_left if hasattr(image_shape, 'crop_left') and image_shape.crop_left else 0.0
    crop_top = image_shape.crop_top if hasattr(image_shape, 'crop_top') and image_shape.crop_top else 0.0
    crop_right = image_shape.crop_right if hasattr(image_shape, 'crop_right') and image_shape.crop_right else 0.0
    crop_bottom = image_shape.crop_bottom if hasattr(image_shape, 'crop_bottom') and image_shape.crop_bottom else 0.0
    return (crop_left, crop_top, crop_right, crop_bottom)


def extract_censor_boxes(slide, presentation) -> list[CensorBox]:
    """
    Extract censorship boxes from a slide.
    Calculates positions RELATIVE TO THE VISIBLE (CROPPED) IMAGE.

    Since we crop the exported image, censor box coordinates are calculated
    relative to what the user will actually see.

    Filters out:
    - Background rectangles (>50% of slide area)
    - Rectangles outside visible image bounds

    Note: This is a heuristic approach. For production use, consider:
    - Naming convention (e.g., shapes named "censor_*")
    - Specific colors or properties
    - Manual tagging in slide notes
    """
    censor_boxes: list[CensorBox] = []
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height
    slide_area = slide_width * slide_height

    # Find the image on this slide
    image_shape = find_image_shape(slide)
    if not image_shape:
        print(f"  Warning: No image found on slide, cannot calculate censor box positions", file=sys.stderr)
        return censor_boxes

    # Image bounds (shape frame in PPTX coordinates)
    img_left = image_shape.left
    img_top = image_shape.top
    img_width = image_shape.width
    img_height = image_shape.height

    # Get crop information
    crop_left, crop_top, crop_right, crop_bottom = get_image_crop_info(image_shape)

    # Calculate the visible (cropped) area within the frame
    # This represents what will actually be in the exported image
    visible_left = img_left + (img_width * crop_left)
    visible_top = img_top + (img_height * crop_top)
    visible_width = img_width * (1.0 - crop_left - crop_right)
    visible_height = img_height * (1.0 - crop_top - crop_bottom)
    visible_right = visible_left + visible_width
    visible_bottom = visible_top + visible_height

    if any([crop_left, crop_top, crop_right, crop_bottom]):
        print(f"  Image cropped: L={crop_left*100:.1f}% T={crop_top*100:.1f}% R={crop_right*100:.1f}% B={crop_bottom*100:.1f}%", file=sys.stderr)
        print(f"  Visible area: ({visible_left:.0f}, {visible_top:.0f}) size: ({visible_width:.0f} × {visible_height:.0f})", file=sys.stderr)
    else:
        print(f"  Image bounds: ({img_left}, {img_top}) size: ({img_width} × {img_height})", file=sys.stderr)

    for shape in slide.shapes:
        # Check if it's an auto shape (rectangle, etc.)
        if shape.shape_type == 1:  # MSO_SHAPE_TYPE.AUTO_SHAPE
            # Check if it has a fill (not an outline-only shape)
            if shape.fill.type == 1:  # SOLID fill
                # Extract position and size
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                # Calculate shape area as percentage of slide
                shape_area = width * height
                area_percentage = (shape_area / slide_area) * 100

                # Skip background rectangles (>50% of slide area)
                if area_percentage > 50:
                    print(f"  Skipping large shape (background): {area_percentage:.1f}% of slide", file=sys.stderr)
                    continue

                # Check if rectangle is within visible (cropped) image bounds
                rect_right = left + width
                rect_bottom = top + height

                within_bounds = (left >= visible_left and
                               rect_right <= visible_right and
                               top >= visible_top and
                               rect_bottom <= visible_bottom)

                if not within_bounds:
                    print(f"  Skipping rectangle outside visible image bounds", file=sys.stderr)
                    continue

                # Calculate position RELATIVE TO VISIBLE (CROPPED) IMAGE
                rel_left = left - visible_left
                rel_top = top - visible_top

                # Convert to percentages of VISIBLE IMAGE dimensions
                x_percent = (rel_left / visible_width) * 100
                y_percent = (rel_top / visible_height) * 100
                width_percent = (width / visible_width) * 100
                height_percent = (height / visible_height) * 100

                # Extract fill color
                try:
                    rgb = shape.fill.fore_color.rgb
                    color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                except Exception:
                    color = "#000000"  # Default to black if color extraction fails

                print(f"  Found censor box: pos=({x_percent:.1f}%, {y_percent:.1f}%), size=({width_percent:.1f}% × {height_percent:.1f}%), color={color}", file=sys.stderr)

                censor_boxes.append(CensorBox(
                    x=round(x_percent, 2),
                    y=round(y_percent, 2),
                    width=round(width_percent, 2),
                    height=round(height_percent, 2),
                    color=color
                ))

    return censor_boxes


def parse_pptx(file_path: Path, category_name: str) -> Category:
    """
    Parse a PPTX file and extract all relevant data.

    Returns a Category object with all slides.
    """
    try:
        prs = Presentation(str(file_path))
    except Exception as e:
        raise ValueError(f"Failed to open PPTX file: {e}")

    slides: list[Slide] = []

    for idx, pptx_slide in enumerate(prs.slides):
        print(f"Processing slide {idx + 1}/{len(prs.slides)}...", file=sys.stderr)

        # Extract image
        image_url = extract_image_as_base64(pptx_slide, idx)
        if not image_url:
            print(f"Warning: No image found on slide {idx + 1}, skipping", file=sys.stderr)
            continue

        # Extract speaker notes (answer)
        answer = extract_speaker_notes(pptx_slide)
        if not answer:
            print(f"Warning: No speaker notes on slide {idx + 1}", file=sys.stderr)

        # Extract censor boxes
        censor_boxes = extract_censor_boxes(pptx_slide, prs)

        slides.append(Slide(
            imageUrl=image_url,
            answer=answer,
            censorBoxes=censor_boxes
        ))

    return Category(name=category_name, slides=slides)


def dataclass_to_dict(obj) -> dict:
    """Convert dataclass to dictionary, handling nested dataclasses."""
    if isinstance(obj, list):
        return [dataclass_to_dict(item) for item in obj]
    elif hasattr(obj, '__dataclass_fields__'):
        return {
            key: dataclass_to_dict(value)
            for key, value in asdict(obj).items()
        }
    else:
        return obj


def main():
    parser = argparse.ArgumentParser(
        description="Parse PPTX file and extract game data as JSON"
    )
    parser.add_argument("input", type=Path, help="Input PPTX file path")
    parser.add_argument("output", type=Path, help="Output JSON file path")
    parser.add_argument(
        "--contestant",
        type=str,
        help="Contestant name (optional, for metadata)"
    )
    parser.add_argument(
        "--category",
        type=str,
        required=True,
        help="Category name (required)"
    )

    args = parser.parse_args()

    # Validate input file
    if not args.input.exists():
        print(f"Error: Input file not found: {args.input}", file=sys.stderr)
        sys.exit(1)

    if not args.input.suffix.lower() == ".pptx":
        print("Error: Input file must be a .pptx file", file=sys.stderr)
        sys.exit(1)

    # Parse PPTX
    print(f"Parsing {args.input}...", file=sys.stderr)
    try:
        category = parse_pptx(args.input, args.category)
    except Exception as e:
        print(f"Error: Failed to parse PPTX: {e}", file=sys.stderr)
        sys.exit(1)

    # Build output data
    metadata = None
    if args.contestant:
        metadata = {"contestantName": args.contestant}

    parsed_data = ParsedData(category=category, metadata=metadata)

    # Convert to dictionary for JSON serialization
    output_dict = dataclass_to_dict(parsed_data)

    # Write output JSON
    print(f"Writing output to {args.output}...", file=sys.stderr)
    try:
        with open(args.output, "w", encoding="utf-8") as f:
            json.dump(output_dict, f, indent=2, ensure_ascii=False)
    except Exception as e:
        print(f"Error: Failed to write output file: {e}", file=sys.stderr)
        sys.exit(1)

    print(f"✓ Successfully parsed {len(category.slides)} slides", file=sys.stderr)
    print(f"✓ Output written to {args.output}", file=sys.stderr)


if __name__ == "__main__":
    main()
