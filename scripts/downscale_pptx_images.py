#!/usr/bin/env python3
"""
Downscale Images in PPTX Files

This script processes PPTX files and downscales embedded images to a maximum
4K resolution (3840x2160) after accounting for any cropping. This reduces the
file size of PPTX files stored in the repository.

The script creates new PPTX files with downscaled images, preserving all other
content including text, shapes, layouts, and crop settings.

Usage:
    python scripts/downscale_pptx_images.py input.pptx output.pptx
    python scripts/downscale_pptx_images.py input_dir/ output_dir/  # Batch mode

Requirements:
    pip install python-pptx pillow
"""

import argparse
import shutil
import sys
from io import BytesIO
from pathlib import Path

try:
    from pptx import Presentation
    from PIL import Image
except ImportError as e:
    print(f"Error: Missing required library: {e}", file=sys.stderr)
    print("Install with: pip install python-pptx pillow", file=sys.stderr)
    sys.exit(1)


def get_image_crop_info(image_shape):
    """
    Get crop information for an image shape.
    Returns (crop_left, crop_top, crop_right, crop_bottom) as fractions (0.0 to 1.0).
    """
    crop_left = image_shape.crop_left if hasattr(image_shape, 'crop_left') and image_shape.crop_left else 0.0
    crop_top = image_shape.crop_top if hasattr(image_shape, 'crop_top') and image_shape.crop_top else 0.0
    crop_right = image_shape.crop_right if hasattr(image_shape, 'crop_right') and image_shape.crop_right else 0.0
    crop_bottom = image_shape.crop_bottom if hasattr(image_shape, 'crop_bottom') and image_shape.crop_bottom else 0.0
    return (crop_left, crop_top, crop_right, crop_bottom)


def calculate_target_size_with_crop(original_width: int, original_height: int,
                                     crop_left: float, crop_top: float,
                                     crop_right: float, crop_bottom: float,
                                     max_width: int = 3840, max_height: int = 2160) -> tuple[int, int]:
    """
    Calculate the target size for downscaling, accounting for crop.

    The goal is to ensure that the VISIBLE (cropped) portion of the image
    is at most 4K resolution.

    Args:
        original_width: Original image width in pixels
        original_height: Original image height in pixels
        crop_left/top/right/bottom: Crop percentages (0.0 to 1.0)
        max_width: Maximum width for visible area (default: 3840)
        max_height: Maximum height for visible area (default: 2160)

    Returns:
        (target_width, target_height) tuple for the full (uncropped) image
    """
    # Calculate visible dimensions after crop
    visible_width = original_width * (1.0 - crop_left - crop_right)
    visible_height = original_height * (1.0 - crop_top - crop_bottom)

    # If visible area is already within limits, no resize needed
    if visible_width <= max_width and visible_height <= max_height:
        return (original_width, original_height)

    # Calculate scale factor based on visible dimensions
    scale = min(max_width / visible_width, max_height / visible_height)

    # Apply scale to FULL image dimensions
    target_width = int(original_width * scale)
    target_height = int(original_height * scale)

    return (target_width, target_height)


def downscale_image(image_bytes: bytes, target_width: int, target_height: int,
                    quality: int = 95) -> bytes:
    """
    Downscale an image to target dimensions.

    Args:
        image_bytes: Original image bytes
        target_width: Target width in pixels
        target_height: Target height in pixels
        quality: JPEG quality (1-100, default: 95 for high quality)

    Returns:
        Downscaled image bytes (as JPEG)
    """
    try:
        img = Image.open(BytesIO(image_bytes))
        original_size = img.size

        # Only resize if target is smaller
        if target_width < original_size[0] or target_height < original_size[1]:
            img = img.resize((target_width, target_height), Image.Resampling.LANCZOS)
            print(f"    Downscaled: {original_size} -> {img.size}", file=sys.stderr)
        else:
            print(f"    No downscale needed: {original_size}", file=sys.stderr)

        # Convert to RGB if needed (handle transparency)
        if img.mode in ('RGBA', 'LA') or (img.mode == 'P' and 'transparency' in img.info):
            background = Image.new('RGB', img.size, (255, 255, 255))
            if img.mode == 'P':
                img = img.convert('RGBA')
            background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
            img = background
        elif img.mode != 'RGB':
            img = img.convert('RGB')

        # Save as JPEG with high quality
        buffer = BytesIO()
        img.save(buffer, format='JPEG', quality=quality, optimize=True)
        return buffer.getvalue()

    except Exception as e:
        print(f"    Warning: Failed to process image: {e}", file=sys.stderr)
        return image_bytes  # Return original on error


def process_pptx(input_path: Path, output_path: Path,
                 max_width: int = 3840, max_height: int = 2160,
                 quality: int = 95) -> bool:
    """
    Process a PPTX file and downscale all images.

    Args:
        input_path: Path to input PPTX file
        output_path: Path to output PPTX file
        max_width: Maximum width for visible area (default: 3840)
        max_height: Maximum height for visible area (default: 2160)
        quality: JPEG quality (default: 95)

    Returns:
        True if successful, False otherwise
    """
    try:
        print(f"Processing: {input_path.name}", file=sys.stderr)
        prs = Presentation(str(input_path))

        total_images = 0
        downscaled_images = 0

        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    total_images += 1
                    print(f"  Slide {slide_idx}: Processing image...", file=sys.stderr)

                    image = shape.image
                    image_bytes = image.blob

                    # Get crop information
                    crop_left, crop_top, crop_right, crop_bottom = get_image_crop_info(shape)

                    # Open image to get dimensions
                    try:
                        img = Image.open(BytesIO(image_bytes))
                        original_width, original_height = img.size

                        # Calculate target size accounting for crop
                        target_width, target_height = calculate_target_size_with_crop(
                            original_width, original_height,
                            crop_left, crop_top, crop_right, crop_bottom,
                            max_width, max_height
                        )

                        # Downscale if needed
                        if target_width < original_width or target_height < original_height:
                            new_image_bytes = downscale_image(
                                image_bytes, target_width, target_height, quality
                            )

                            # Replace image in shape
                            # Note: python-pptx doesn't have a direct API to replace image bytes,
                            # so we need to work with the underlying XML
                            image_part = shape.image.image_part
                            image_part._blob = new_image_bytes

                            downscaled_images += 1
                        else:
                            print(f"    Visible area already within {max_width}x{max_height}", file=sys.stderr)

                    except Exception as e:
                        print(f"    Warning: Failed to process image: {e}", file=sys.stderr)

        # Save the modified presentation
        prs.save(str(output_path))

        print(f"✓ Processed {total_images} images ({downscaled_images} downscaled)", file=sys.stderr)
        print(f"✓ Saved: {output_path}", file=sys.stderr)

        # Show file size comparison
        input_size = input_path.stat().st_size / (1024 * 1024)  # MB
        output_size = output_path.stat().st_size / (1024 * 1024)  # MB
        savings = ((input_size - output_size) / input_size * 100) if input_size > 0 else 0
        print(f"  File size: {input_size:.2f}MB -> {output_size:.2f}MB ({savings:.1f}% reduction)", file=sys.stderr)

        return True

    except Exception as e:
        print(f"✗ Failed to process {input_path.name}: {e}", file=sys.stderr)
        return False


def main():
    parser = argparse.ArgumentParser(
        description="Downscale images in PPTX files to 4K resolution"
    )
    parser.add_argument("input", type=Path, help="Input PPTX file or directory")
    parser.add_argument("output", type=Path, help="Output PPTX file or directory")
    parser.add_argument(
        "--max-width",
        type=int,
        default=3840,
        help="Maximum width for visible area (default: 3840)"
    )
    parser.add_argument(
        "--max-height",
        type=int,
        default=2160,
        help="Maximum height for visible area (default: 2160)"
    )
    parser.add_argument(
        "--quality",
        type=int,
        default=95,
        help="JPEG quality 1-100 (default: 95)"
    )

    args = parser.parse_args()

    # Validate input
    if not args.input.exists():
        print(f"Error: Input not found: {args.input}", file=sys.stderr)
        sys.exit(1)

    # Single file mode
    if args.input.is_file():
        if not args.input.suffix.lower() == ".pptx":
            print("Error: Input file must be a .pptx file", file=sys.stderr)
            sys.exit(1)

        # Create output directory if needed
        args.output.parent.mkdir(parents=True, exist_ok=True)

        success = process_pptx(args.input, args.output, args.max_width, args.max_height, args.quality)
        sys.exit(0 if success else 1)

    # Batch directory mode
    elif args.input.is_dir():
        args.output.mkdir(parents=True, exist_ok=True)

        pptx_files = list(args.input.glob("*.pptx"))
        if not pptx_files:
            print(f"No PPTX files found in: {args.input}", file=sys.stderr)
            sys.exit(1)

        print(f"\nFound {len(pptx_files)} PPTX files\n", file=sys.stderr)

        successful = 0
        failed = 0

        for pptx_file in pptx_files:
            output_file = args.output / pptx_file.name
            if process_pptx(pptx_file, output_file, args.max_width, args.max_height, args.quality):
                successful += 1
            else:
                failed += 1
            print()  # Blank line between files

        print("=" * 50, file=sys.stderr)
        print(f"Complete: {successful} successful, {failed} failed", file=sys.stderr)
        sys.exit(1 if failed > 0 else 0)

    else:
        print(f"Error: Invalid input: {args.input}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
