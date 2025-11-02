#!/usr/bin/env python3
"""
Diagnostic script to examine PPTX structure and coordinate systems.
This helps understand how images and shapes are positioned.
"""

import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError as e:
    print(f"Error: Missing required library: {e}", file=sys.stderr)
    print("Install with: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


def diagnose_slide(prs, slide_index):
    """Examine a single slide in detail."""
    slide = prs.slides[slide_index]

    print(f"\n{'='*80}")
    print(f"SLIDE {slide_index + 1} DIAGNOSIS")
    print(f"{'='*80}")

    # Slide dimensions
    print(f"\nSlide Dimensions:")
    print(f"  Width:  {prs.slide_width} EMU ({prs.slide_width / 914400:.2f} inches)")
    print(f"  Height: {prs.slide_height} EMU ({prs.slide_height / 914400:.2f} inches)")

    # Examine all shapes
    print(f"\nShapes on slide ({len(slide.shapes)} total):")
    print(f"{'-'*80}")

    image_shapes = []
    rect_shapes = []

    for shape_idx, shape in enumerate(slide.shapes):
        print(f"\nShape {shape_idx + 1}:")
        print(f"  Name: {shape.name}")
        print(f"  Type: {shape.shape_type} ({get_shape_type_name(shape.shape_type)})")

        # Position and size
        print(f"  Position (EMU):")
        print(f"    Left: {shape.left} ({shape.left / 914400:.2f} inches)")
        print(f"    Top:  {shape.top} ({shape.top / 914400:.2f} inches)")
        print(f"  Size (EMU):")
        print(f"    Width:  {shape.width} ({shape.width / 914400:.2f} inches)")
        print(f"    Height: {shape.height} ({shape.height / 914400:.2f} inches)")

        # Calculate percentages relative to slide
        x_percent = (shape.left / prs.slide_width) * 100
        y_percent = (shape.top / prs.slide_height) * 100
        width_percent = (shape.width / prs.slide_width) * 100
        height_percent = (shape.height / prs.slide_height) * 100

        print(f"  Position (% of slide):")
        print(f"    X: {x_percent:.2f}%")
        print(f"    Y: {y_percent:.2f}%")
        print(f"  Size (% of slide):")
        print(f"    Width:  {width_percent:.2f}%")
        print(f"    Height: {height_percent:.2f}%")

        # Special handling for pictures
        if shape.shape_type == 13:  # PICTURE
            image_shapes.append((shape_idx, shape))
            print(f"  📷 IMAGE DETECTED")
            if hasattr(shape, 'image'):
                print(f"    Image format: {shape.image.ext}")
                print(f"    Image size: {len(shape.image.blob)} bytes")

            # Check for cropping
            if hasattr(shape, 'crop_left'):
                crop_left = shape.crop_left if shape.crop_left else 0
                crop_top = shape.crop_top if shape.crop_top else 0
                crop_right = shape.crop_right if shape.crop_right else 0
                crop_bottom = shape.crop_bottom if shape.crop_bottom else 0

                if any([crop_left, crop_top, crop_right, crop_bottom]):
                    print(f"  ✂️  IMAGE IS CROPPED:")
                    print(f"    Crop left:   {crop_left:.4f} ({crop_left * 100:.2f}%)")
                    print(f"    Crop top:    {crop_top:.4f} ({crop_top * 100:.2f}%)")
                    print(f"    Crop right:  {crop_right:.4f} ({crop_right * 100:.2f}%)")
                    print(f"    Crop bottom: {crop_bottom:.4f} ({crop_bottom * 100:.2f}%)")

                    # Calculate visible area
                    visible_width = 1.0 - crop_left - crop_right
                    visible_height = 1.0 - crop_top - crop_bottom
                    print(f"    Visible area: {visible_width * 100:.2f}% width × {visible_height * 100:.2f}% height")
                else:
                    print(f"    No cropping applied")

        # Special handling for autoshapes (rectangles)
        if shape.shape_type == 1:  # AUTO_SHAPE
            if hasattr(shape, 'fill') and shape.fill.type == 1:  # SOLID fill
                rect_shapes.append((shape_idx, shape))
                try:
                    rgb = shape.fill.fore_color.rgb
                    color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                    print(f"  🟦 FILLED RECTANGLE")
                    print(f"    Color: {color}")
                except:
                    print(f"  🟦 FILLED RECTANGLE (color extraction failed)")

    # Analysis
    print(f"\n{'='*80}")
    print(f"ANALYSIS:")
    print(f"{'='*80}")
    print(f"Images found: {len(image_shapes)}")
    print(f"Filled rectangles found: {len(rect_shapes)}")

    if image_shapes and rect_shapes:
        print(f"\n🔍 COORDINATE RELATIONSHIP ANALYSIS:")
        for img_idx, img_shape in image_shapes:
            print(f"\n  Image (Shape {img_idx + 1}):")
            print(f"    Bounds: ({img_shape.left}, {img_shape.top}) to "
                  f"({img_shape.left + img_shape.width}, {img_shape.top + img_shape.height})")

            for rect_idx, rect_shape in rect_shapes:
                print(f"\n  Rectangle (Shape {rect_idx + 1}) relative to image:")

                # Calculate position relative to image
                rel_left = rect_shape.left - img_shape.left
                rel_top = rect_shape.top - img_shape.top

                print(f"    Absolute position: ({rect_shape.left}, {rect_shape.top})")
                print(f"    Relative to image: ({rel_left}, {rel_top})")

                # Calculate percentages relative to image
                x_pct_of_image = (rel_left / img_shape.width) * 100
                y_pct_of_image = (rel_top / img_shape.height) * 100
                width_pct_of_image = (rect_shape.width / img_shape.width) * 100
                height_pct_of_image = (rect_shape.height / img_shape.height) * 100

                print(f"    Position (% of IMAGE):")
                print(f"      X: {x_pct_of_image:.2f}%")
                print(f"      Y: {y_pct_of_image:.2f}%")
                print(f"    Size (% of IMAGE):")
                print(f"      Width:  {width_pct_of_image:.2f}%")
                print(f"      Height: {height_pct_of_image:.2f}%")

                # Check if rectangle is within image bounds
                rect_right = rect_shape.left + rect_shape.width
                rect_bottom = rect_shape.top + rect_shape.height
                img_right = img_shape.left + img_shape.width
                img_bottom = img_shape.top + img_shape.height

                within_bounds = (rect_shape.left >= img_shape.left and
                               rect_right <= img_right and
                               rect_shape.top >= img_shape.top and
                               rect_bottom <= img_bottom)

                if within_bounds:
                    print(f"    ✅ Rectangle is WITHIN image bounds")
                else:
                    print(f"    ⚠️  Rectangle is OUTSIDE/OVERLAPS image bounds")


def get_shape_type_name(shape_type):
    """Get human-readable shape type name."""
    types = {
        1: "AUTO_SHAPE",
        13: "PICTURE",
        14: "PLACEHOLDER",
        17: "TEXT_BOX",
    }
    return types.get(shape_type, f"UNKNOWN_{shape_type}")


def main():
    if len(sys.argv) < 2:
        print("Usage: python diagnose_pptx.py <file.pptx> [slide_number]")
        print("\nExample:")
        print("  python diagnose_pptx.py presentation.pptx 2")
        print("  (diagnoses slide 2, or all slides if number not specified)")
        sys.exit(1)

    file_path = Path(sys.argv[1])

    if not file_path.exists():
        print(f"Error: File not found: {file_path}", file=sys.stderr)
        sys.exit(1)

    try:
        prs = Presentation(str(file_path))
    except Exception as e:
        print(f"Error: Failed to open PPTX: {e}", file=sys.stderr)
        sys.exit(1)

    print(f"Opened: {file_path}")
    print(f"Total slides: {len(prs.slides)}")

    # Determine which slides to diagnose
    if len(sys.argv) >= 3:
        slide_num = int(sys.argv[2])
        if slide_num < 1 or slide_num > len(prs.slides):
            print(f"Error: Slide number must be between 1 and {len(prs.slides)}", file=sys.stderr)
            sys.exit(1)
        diagnose_slide(prs, slide_num - 1)
    else:
        # Diagnose all slides
        for i in range(len(prs.slides)):
            diagnose_slide(prs, i)


if __name__ == "__main__":
    main()
